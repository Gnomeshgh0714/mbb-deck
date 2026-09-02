#!/usr/bin/env python3
"""
mbb-deck · PPTX 路径生成器（generate_pptx.py）

把 plan.json 渲染为原生可编辑的 .pptx（16:9）。工作路径参考 lampertb / mbb-decks /
mckinsey-pptx 三套 PPTX skill 的共识流程：

    plan.json（Step 1-4 语义层 + Step 5B 内容级扩展）
      -> 本脚本确定性渲染
      -> 输出机械质量告警
      -> 发现问题改 plan.json 重跑（不手改 .pptx）

框架版状态：
  [已实现] plan.json 结构校验、python-pptx 依赖自检（缺失给安装提示）、
           16:9 画布、每页结构骨架（标题带 action title + 页脚带来源注/页码）、
           质量告警框架
  [🔲 待填] 版式内容渲染（依赖 pptx-generation.md §3 映射表定稿）、
           主题预设映射（design-system.md §1 三套预设的 pptx 色值）、
           质量告警规则（pptx-generation.md §4）

用法：
  python3 generate_pptx.py --plan plan.json --output deck.pptx

退出码：0 = 生成成功（可有 warning）；1 = plan.json 校验失败或渲染错误；2 = 依赖缺失。
"""

import argparse
import json
import sys
from pathlib import Path

# ---------------------------------------------------------------------------
# 契约（与 references/routing.md §4.2 同步）
# ---------------------------------------------------------------------------

PAGE_REQUIRED_FIELDS = ("page_no", "role", "action_title", "layout", "shape", "not", "fit_check")
PLAN_REQUIRED_FIELDS = ("format", "governing_thought", "spine", "pages")

# 版式 ID -> 渲染器注册表。
# 🔲 待填 -- 定稿后按 references/pptx-generation.md §3 映射表逐版式补齐渲染函数；
# 当前仅注册结构级通用渲染（封面 + 标题带/页脚带骨架）。
LAYOUT_RENDERERS = {}


# ---------------------------------------------------------------------------
# 校验
# ---------------------------------------------------------------------------

def validate_plan(plan: dict) -> tuple[list[str], list[str]]:
    """返回 (blockers, warnings)。blockers 阻止渲染，warnings 照常生成但须披露。"""
    blockers, warnings = [], []

    for field in PLAN_REQUIRED_FIELDS:
        if field not in plan:
            blockers.append(f"plan.json 缺少必填字段：{field}")

    if plan.get("format") != "pptx":
        blockers.append(f'format 字段应为 "pptx"，实际为 {plan.get("format")!r}（HTML 路径不经过本脚本）')

    pages = plan.get("pages") or []
    if not isinstance(pages, list) or not pages:
        blockers.append("pages 为空或不是列表")
        return blockers, warnings

    for i, page in enumerate(pages, 1):
        missing = [f for f in PAGE_REQUIRED_FIELDS if f not in page]
        if missing:
            blockers.append(f"第 {i} 页缺少字段：{', '.join(missing)}")
        if not page.get("action_title"):
            warnings.append(f"第 {i} 页 action_title 为空（铁律 1：标题必须是结论句）")
        if "content" not in page:
            # 框架版：内容级契约未定稿，content 缺失时仅渲染结构骨架
            warnings.append(f"第 {i} 页无 content 块：仅渲染结构骨架（内容级契约见 pptx-generation.md §2，🔲 待填）")

    return blockers, warnings


# ---------------------------------------------------------------------------
# 依赖自检（mbb-decks 模式）
# ---------------------------------------------------------------------------

def import_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        return Presentation, (Inches, Pt)
    except ImportError:
        print("错误：渲染器需要 python-pptx。请在终端运行：", file=sys.stderr)
        print("  pip install python-pptx", file=sys.stderr)
        print("装好后重新执行本脚本。", file=sys.stderr)
        return None, None


# ---------------------------------------------------------------------------
# 结构骨架渲染（版式无关部分，对应 assets/deck-template.html 的已定结构）
# ---------------------------------------------------------------------------

def render_skeleton(prs, slide, page: dict, page_no: int, total: int):
    """渲染标题带 + 页脚带 + 内容占位区。字号/色值为中性占位，🔲 待填：
    定稿后按 design-system.md §1/§3 替换为预设令牌。"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    INK = RGBColor(0x11, 0x11, 0x11)      # 🔲 待填：预设 ink
    NAVY = RGBColor(0x11, 0x11, 0x11)     # 🔲 待填：预设 navy
    MUTED = RGBColor(0x66, 0x66, 0x66)    # 🔲 待填：预设 muted
    LINE = RGBColor(0xDD, 0xDD, 0xDD)     # 🔲 待填：预设 line

    is_cover = page.get("layout") == "cover"

    if is_cover:
        # 封面：居中主标题（结构对应 deck-template.html 的 .slide.cover）
        box = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.3), Inches(1.9))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = page["action_title"]
        p.font.size = Pt(40)              # 🔲 待填：封面标题字级
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER
        return

    # 标题带：action title（结论句）
    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.9))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = page["action_title"]
    p.font.size = Pt(24)                  # 🔲 待填：action title 字级
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 标题带下沿细分隔线（🔲 待填：粗细 0.5-1pt 定稿后确认）
    from pptx.enum.shapes import MSO_SHAPE
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.45), Inches(11.8), Pt(1))
    rule.fill.solid()
    rule.fill.fore_color.rgb = LINE
    rule.line.fill.background()

    # 页脚带：左下来源注 + 右下页码
    source = page.get("source")
    if source:
        src_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.9), Inches(9.0), Inches(0.4))
        sp = src_box.text_frame.paragraphs[0]
        sp.text = source
        sp.font.size = Pt(8)              # 🔲 待填：来源注字级（8pt 斜体）
        sp.font.italic = True
        sp.font.color.rgb = MUTED

    no_box = slide.shapes.add_textbox(Inches(12.0), Inches(6.9), Inches(0.55), Inches(0.4))
    np_ = no_box.text_frame.paragraphs[0]
    np_.text = f"{page_no} / {total}"
    np_.font.size = Pt(8)                 # 🔲 待填：页码字级
    np_.font.color.rgb = MUTED
    np_.alignment = PP_ALIGN.RIGHT


# ---------------------------------------------------------------------------
# 入口
# ---------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(description="mbb-deck · PPTX 路径生成器")
    parser.add_argument("--plan", type=Path, required=True, help="内容级 plan.json")
    parser.add_argument("--output", type=Path, required=True, help="输出 .pptx 路径")
    args = parser.parse_args()

    if not args.plan.is_file():
        print(f"错误：plan 文件不存在：{args.plan}", file=sys.stderr)
        return 1
    try:
        plan = json.loads(args.plan.read_text(encoding="utf-8"))
    except json.JSONDecodeError as e:
        print(f"错误：plan.json 不是合法 JSON：{e}", file=sys.stderr)
        return 1

    blockers, warnings = validate_plan(plan)
    if blockers:
        print("plan.json 校验未通过（blocker）：")
        for b in blockers:
            print(f"  [blocker] {b}")
        return 1

    Presentation, _ = import_pptx()
    if Presentation is None:
        return 2

    prs = Presentation()
    prs.slide_width = 13214400    # 13.333 in（EMU）· 16:9
    prs.slide_height = 6858000    # 7.5 in（EMU）
    blank = prs.slide_layouts[6]  # 空白版式，全部对象自绘

    pages = plan["pages"]
    total = len(pages)
    for page_no, page in enumerate(pages, 1):
        slide = prs.slides.add_slide(blank)
        render_skeleton(prs, slide, page, page_no, total)

        # 版式内容渲染：🔲 待填（定稿后按 LAYOUT_RENDERERS 注册表分发）
        layout = page.get("layout", "?")
        renderer = LAYOUT_RENDERERS.get(layout)
        if renderer:
            renderer(prs, slide, page)
        elif page.get("content"):
            warnings.append(
                f"第 {page_no} 页版式 {layout} 无渲染器（映射表待定稿，pptx-generation.md §3 🔲 待填）：content 未渲染"
            )

    prs.save(str(args.output))

    # 质量告警汇总（🔲 待填：定稿后扩充为 pptx-generation.md §4 的完整清单）
    print(f"已生成：{args.output}（{total} 页 · 16:9 · 原生可编辑对象）")
    if warnings:
        print(f"质量告警：{len(warnings)} 条（修复方式：改 plan.json 后重跑，不手改 .pptx）")
        for w in warnings:
            print(f"  [warning] {w}")
    else:
        print("机械质量告警：无")
    return 0


if __name__ == "__main__":
    sys.exit(main())
